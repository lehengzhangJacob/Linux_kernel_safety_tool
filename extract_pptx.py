from pptx import Presentation
import sys

try:
    prs = Presentation("演示文稿1.pptx")
    print(f"Presentation loaded. Total slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
except ImportError:
    print("Error: python-pptx is not installed.")
except Exception as e:
    print(f"Error processing file: {e}")
